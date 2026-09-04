from __future__ import annotations

import json
import os
import re
import subprocess
import sys
import zipfile
from html import escape
from io import BytesIO
from pathlib import Path
from typing import Any, cast
from xml.etree import ElementTree

import pytest
from pptx import Presentation

from archsift.canonical import JsonObject
from archsift.executive_summary import (
    PART_TITLES,
    ExecutiveSummary,
    SummaryPart,
    SummaryStatement,
    build_executive_summary,
)
from archsift.html_report import render_executive_html_report
from archsift.masking import MASKING_WARNING
from archsift.pptx_report import (
    POINTS_PER_SLIDE,
    render_executive_pptx_report,
    render_executive_summary_pptx,
)
from archsift.record_view import ReportRecordError
from archsift.report_text import visible_text

_GOLDEN_DIR = Path(__file__).parent / "golden"
_POSITIVE_RECORD = _GOLDEN_DIR / "decision-record-positive-v1.json"
_ABSTENTION_RECORD = _GOLDEN_DIR / "decision-record-abstention-veto-v1.json"
_HTML_GOLDEN = _GOLDEN_DIR / "executive-summary-abstention-veto-v1.html"
_PPTX_GOLDEN = _GOLDEN_DIR / "executive-summary-abstention-veto-v1.pptx"

# Payloads that would escape a text node, an attribute, a URL, a raw-text
# element, or a DrawingML text run if authored strings were ever emitted as
# markup rather than as data.
_INJECTION_PAYLOAD = (
    '<script>alert(1)</script><img src=x onerror="alert(2)">'
    "</textarea><textarea>forged</textarea>"
    "</style><style>body{display:none}</style>"
    '</title><a href="javascript:alert(3)">forged link</a>'
    "</a:t></a:r></a:p><a:p><a:r><a:t>forged slide text</a:t></a:r></a:p>"
    "</dd></dl><h1>forged heading</h1>"
    "\"'&<>` \x00\x1b\u200b\u2028\u202e caf\u00e9"
)

_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
_SLIDE_PART = re.compile(r"ppt/slides/slide(\d+)\.xml")


def _record(path: Path = _ABSTENTION_RECORD) -> JsonObject:
    return cast(JsonObject, json.loads(path.read_bytes()))


def _slide_index(name: str) -> int:
    match = _SLIDE_PART.fullmatch(name)
    assert match is not None
    return int(match.group(1))


def _slide_text(deck: bytes) -> list[list[str]]:
    """Return each slide's text runs, read back from the written package."""
    slides: list[list[str]] = []
    with zipfile.ZipFile(BytesIO(deck)) as archive:
        names = sorted(
            (name for name in archive.namelist() if _SLIDE_PART.fullmatch(name)),
            key=_slide_index,
        )
        for name in names:
            root = ElementTree.fromstring(archive.read(name))
            slides.append([node.text or "" for node in root.iter(f"{_A}t")])
    return slides


def _summary_lines(record: JsonObject) -> list[str]:
    summary = build_executive_summary(record)
    return [
        f"{statement.label}: {statement.text}"
        for part in summary.parts
        for statement in part.statements
    ]


def test_executive_html_matches_its_exact_golden() -> None:
    content = render_executive_html_report(_record())

    assert content == _HTML_GOLDEN.read_bytes()
    assert content.endswith(b"\n") and not content.endswith(b"\n\n")
    assert b"\r" not in content
    text = content.decode("utf-8")
    assert "<title>ArchSift Executive Summary</title>" in text
    assert re.findall(r"<h2>(.*?)</h2>", text) == list(PART_TITLES)
    assert "More evidence is needed before an option can be indicated." in text
    assert "The fictional disposition would be released without approval." in text
    assert "<dt>Masking notice</dt>" in text


def test_executive_pptx_matches_its_exact_golden() -> None:
    deck = render_executive_pptx_report(_record())

    assert deck == _PPTX_GOLDEN.read_bytes()
    assert deck[:2] == b"PK"


def test_pptx_golden_is_pinned_as_binary_content() -> None:
    attributes = (Path(__file__).parent.parent / ".gitattributes").read_text(encoding="utf-8")

    assert "tests/golden/*.pptx binary" in attributes


def test_both_executive_formats_state_the_same_facts() -> None:
    """FR-017: HTML and PPTX render one summary, so neither can drift."""
    record = _record()
    html = render_executive_html_report(record).decode("utf-8")
    slides = _slide_text(render_executive_pptx_report(record))
    deck_text = "\n".join(run for slide in slides for run in slide)

    lines = _summary_lines(record)
    assert len(lines) >= 15
    for line in lines:
        label, _, values = line.partition(": ")
        assert f"<dt>{escape(visible_text(label), quote=True)}</dt>" in html, label
        assert escape(visible_text(values), quote=True) in html, line
        assert f"{label}: {visible_text(values)}" in deck_text, line
    for title in PART_TITLES:
        assert f"<h2>{title}</h2>" in html
        assert title in deck_text


def test_pptx_package_is_a_complete_offline_presentation() -> None:
    deck = render_executive_pptx_report(_record())

    with zipfile.ZipFile(BytesIO(deck)) as archive:
        names = archive.namelist()
        assert archive.testzip() is None
        for required in (
            "[Content_Types].xml",
            "_rels/.rels",
            "ppt/presentation.xml",
            "ppt/_rels/presentation.xml.rels",
            "ppt/slideMasters/slideMaster1.xml",
            "ppt/slideLayouts/slideLayout1.xml",
            "ppt/theme/theme1.xml",
            "docProps/core.xml",
            "docProps/app.xml",
        ):
            assert required in names, required
        blob = b"".join(archive.read(name) for name in names)
        for name in names:
            ElementTree.fromstring(archive.read(name))
    # No media part, and no reference that would leave the package.
    assert not [name for name in names if name.startswith("ppt/media/")]
    assert b"http://schemas.openxmlformats.org" in blob
    external = re.findall(rb"https?://(?!schemas\.openxmlformats\.org|purl\.org)[^\s\"'<>]+", blob)
    assert external == [], external
    assert b"<a:blip" not in blob and b"embed=" not in blob
    assert b'TargetMode="External"' not in blob


def test_pptx_bytes_are_reproducible_and_carry_no_timestamp() -> None:
    record = _record()

    first = render_executive_pptx_report(record)
    second = render_executive_pptx_report(record)

    assert first == second
    with zipfile.ZipFile(BytesIO(first)) as archive:
        infos = archive.infolist()
        assert {info.date_time for info in infos} == {(1980, 1, 1, 0, 0, 0)}
        assert {info.compress_type for info in infos} == {zipfile.ZIP_STORED}
        # `create_system` would otherwise record Windows or POSIX.
        assert {info.create_system for info in infos} == {0}
        assert {info.external_attr for info in infos} == {0o600 << 16}
        assert {info.flag_bits for info in infos} == {0}
        assert {info.extract_version for info in infos} == {20}
        core = archive.read("docProps/core.xml").decode("utf-8")
    assert "dcterms:created" not in core and "dcterms:modified" not in core


def test_executive_outputs_are_hash_seed_independent() -> None:
    script = f"""
import json
from pathlib import Path
from archsift.html_report import render_executive_html_report
from archsift.pptx_report import render_executive_pptx_report
record = json.loads(Path({str(_ABSTENTION_RECORD)!r}).read_bytes())
print(render_executive_html_report(record).hex())
print(render_executive_pptx_report(record).hex())
"""
    outputs = []
    for seed in ("1", "947"):
        environment = os.environ.copy()
        environment["PYTHONHASHSEED"] = seed
        outputs.append(
            subprocess.run(
                [sys.executable, "-c", script],
                check=True,
                capture_output=True,
                text=True,
                env=environment,
            ).stdout
        )
    assert outputs[0] == outputs[1]


def test_authored_markup_and_xml_render_as_inert_text_in_both_formats() -> None:
    record = _record(_POSITIVE_RECORD)
    dossier = cast(dict[str, Any], record["dossier"])
    dossier["case"]["title"] = _INJECTION_PAYLOAD
    dossier["autonomy_permission"]["hard_vetoes"][0]["condition"] = _INJECTION_PAYLOAD

    html = render_executive_html_report(record).decode("utf-8")
    deck = render_executive_pptx_report(record)

    assert html.count("<script") == 0 and html.count("<a ") == 0
    assert html.count("<style>") == 1 and html.count("<title>") == 1
    assert set(re.findall(r'([a-zA-Z-]+)="', html)) == {
        "charset",
        "class",
        "content",
        "lang",
        "name",
    }
    assert escape(visible_text(_INJECTION_PAYLOAD), quote=True) in html

    slides = _slide_text(deck)
    baseline = _slide_text(render_executive_pptx_report(_record(_POSITIVE_RECORD)))
    # The payload stayed inside single text runs: it created no new run, no new
    # paragraph, and no new slide.
    assert len(slides) == len(baseline)
    assert [len(slide) for slide in slides] == [len(slide) for slide in baseline]
    assert any(visible_text(_INJECTION_PAYLOAD) in run for slide in slides for run in slide)
    assert not any(run.strip() == "forged slide text" for slide in slides for run in slide)
    with zipfile.ZipFile(BytesIO(deck)) as archive:
        for name in archive.namelist():
            ElementTree.fromstring(archive.read(name))
        blob = b"".join(archive.read(name) for name in archive.namelist())
    assert b"<script" not in blob
    assert b"\x00" not in blob and b"\x1b" not in blob


def test_masking_reaches_both_executive_formats() -> None:
    record = _record(_POSITIVE_RECORD)
    dossier = cast(dict[str, Any], record["dossier"])
    dossier["case"]["title"] = "Card 4111 1111 1111 1111 and api_key: AKIAIOSFODNN7EXAMPLE"
    record.pop("masking", None)

    html = render_executive_html_report(record).decode("utf-8")
    deck_text = "\n".join(
        run for slide in _slide_text(render_executive_pptx_report(record)) for run in slide
    )

    for rendered in (html, deck_text):
        assert "4111 1111 1111 1111" not in rendered
        assert "AKIAIOSFODNN7EXAMPLE" not in rendered
        assert "[ARCHSIFT-MASKED:payment-card]" in rendered
        assert "[ARCHSIFT-MASKED:credential]" in rendered
    assert MASKING_WARNING in html and MASKING_WARNING in deck_text


def test_both_formats_restate_the_record_identity_and_derive_no_other() -> None:
    record = _record()
    identity = cast(str, record["record_content_identity"])
    html = render_executive_html_report(record).decode("utf-8")
    deck_text = "\n".join(
        run for slide in _slide_text(render_executive_pptx_report(record)) for run in slide
    )

    assert identity in html and identity in deck_text
    assert set(re.findall(r"sha256:[0-9a-f]{64}", html)) == {identity}
    assert set(re.findall(r"sha256:[0-9a-f]{64}", deck_text)) == {identity}


def test_deck_frames_every_part_with_a_title_and_a_masking_notice() -> None:
    record = _record()
    slides = _slide_text(render_executive_pptx_report(record))
    titles = [slide[0] for slide in slides]

    assert titles[0] == "ArchSift Executive Summary"
    assert titles[-1] == "Masking Notice"
    assert [title for title in titles if title in PART_TITLES] == list(PART_TITLES)
    assert all(title in PART_TITLES or title.endswith(" (continued)") for title in titles[1:-1]), (
        titles
    )
    for slide in slides:
        # One title run plus at most a full page of bullet runs.
        assert len(slide) <= POINTS_PER_SLIDE + 1
    rendered = {run for slide in slides for run in slide}
    for line in _summary_lines(record):
        label, _, values = line.partition(": ")
        assert f"{label}: {visible_text(values)}" in rendered, line


def test_an_oversized_part_continues_onto_further_slides_without_truncation() -> None:
    statements = tuple(
        SummaryStatement("Option", f"Synthetic option {index}. No flag is raised on this option.")
        for index in range(POINTS_PER_SLIDE * 2 + 1)
    )
    summary = ExecutiveSummary(
        record_content_identity="sha256:" + "0" * 64,
        vocabulary_version="0.0.0-test",
        case_title="Oversized synthetic part",
        language="en",
        parts=(
            SummaryPart("Summary", (SummaryStatement("The task", "Synthetic."),)),
            SummaryPart("Business analysis", (SummaryStatement("Step 1", "Synthetic."),)),
            SummaryPart("Result and reasoning", statements),
        ),
    )

    slides = _slide_text(render_executive_summary_pptx(summary))
    titles = [slide[0] for slide in slides]

    assert titles == [
        "ArchSift Executive Summary",
        "Summary",
        "Business analysis",
        "Result and reasoning",
        "Result and reasoning (continued)",
        "Result and reasoning (continued)",
        "Masking Notice",
    ]
    bullets = [run for slide in slides[3:6] for run in slide[1:]]
    assert len(bullets) == len(statements)
    for statement in statements:
        assert f"{statement.label}: {statement.text}" in bullets


def test_executive_renderers_fail_closed_on_an_unusable_record() -> None:
    record = _record()
    record.pop("assessment")

    with pytest.raises(ReportRecordError, match=r"\$ is missing assessment"):
        render_executive_html_report(record)
    with pytest.raises(ReportRecordError, match=r"\$ is missing assessment"):
        render_executive_pptx_report(record)


def test_deck_opens_as_a_real_presentation_under_an_independent_reader() -> None:
    """The written package is a presentation, not merely a well-formed zip.

    `python-pptx` is a development-only dependency: it never reaches the
    shipped package, and it reads the deck here rather than writing it, so the
    deterministic writer is checked against an independent implementation of
    the format.
    """
    record = _record()

    presentation = Presentation(BytesIO(render_executive_pptx_report(record)))

    slides = list(presentation.slides)
    assert presentation.slide_width == 12192000
    assert presentation.slide_height == 6858000
    assert len(slides) == len(_slide_text(render_executive_pptx_report(record)))
    read_back = [
        "\n".join(shape.text_frame.text for shape in slide.shapes if shape.has_text_frame)
        for slide in slides
    ]
    assert read_back[0].startswith("ArchSift Executive Summary")
    assert cast(str, record["record_content_identity"]) in read_back[0]
    document = "\n".join(read_back)
    for line in _summary_lines(record):
        label, _, values = line.partition(": ")
        assert f"{label}: {visible_text(values)}" in document, line
    assert MASKING_WARNING in document


def test_the_shipped_package_never_imports_the_verification_reader() -> None:
    """python-pptx is a development dependency, never a runtime one."""
    source = Path(__file__).parent.parent / "src" / "archsift"
    modules = sorted(source.glob("*.py"))

    assert len(modules) >= 15
    for module in modules:
        text = module.read_text(encoding="utf-8")
        assert "import pptx" not in text, module.name
        assert "from pptx" not in text, module.name

    installed = subprocess.run(
        [sys.executable, "-c", "import sys, archsift.pptx_report; print('pptx' in sys.modules)"],
        check=True,
        capture_output=True,
        text=True,
    )
    assert installed.stdout.strip() == "False"


def test_rendering_a_deck_pulls_no_network_stack_into_the_import_graph() -> None:
    """NFR-001: writing a presentation must not import a networking module.

    `xml.sax.saxutils` reaches `urllib.request` and `ssl`; escaping the three
    markup characters locally keeps them out of a tool that never opens a
    connection.

    The check is differential rather than absolute. Some supported Python
    versions already import `socket` for the package's own version lookup
    (`importlib.metadata` reaches `email.utils`), so what matters is that the
    report modules add nothing to that baseline.
    """
    probe = (
        "import sys\n"
        "import archsift\n"
        "baseline = set(sys.modules)\n"
        "import archsift.cli\n"
        "from archsift.pptx_report import render_executive_pptx_report\n"
        "from archsift.html_report import render_executive_html_report\n"
        "network = {'ssl', 'socket', 'http.client', 'urllib.request', 'asyncio',\n"
        "           'ftplib', 'smtplib', 'telnetlib', 'webbrowser', 'xmlrpc.client'}\n"
        "print(sorted((set(sys.modules) - baseline) & network))\n"
    )

    imported = subprocess.run(
        [sys.executable, "-c", probe], check=True, capture_output=True, text=True
    )

    assert imported.stdout.strip() == "[]", imported.stdout
